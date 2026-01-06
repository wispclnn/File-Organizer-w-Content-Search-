import os
import sys
import json
import shutil
import time
from datetime import datetime
from pathlib import Path
import threading  # --- Content Search Added ---
import traceback  # --- Content Search Added ---
import re  # --- Content Search Added ---
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


import tkinter as tk
from tkinter import filedialog, messagebox

# We gracefully degrade if some libraries are missing.
try:
    import fitz  
except Exception:
    fitz = None

try:
    import docx  # python-docx for .docx
except Exception: 
    docx = None

try:
    from pptx import Presentation  # python-pptx for .pptx
except Exception: 
    Presentation = None

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:  
    TfidfVectorizer = None
    cosine_similarity = None


try:
    from PIL import Image, ImageTk
except Exception:  
    Image = None
    ImageTk = None


try:
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except Exception:
    pytesseract = None


if getattr(sys, 'frozen', False):
    APP_DIR = os.path.dirname(sys.executable)
else:
    APP_DIR = os.path.dirname(os.path.abspath(__file__))
    
RULES_PATH = os.path.join(APP_DIR, 'rules.json')
TAGS_PATH = os.path.join(APP_DIR, 'tags.json')
FOLDERS_PATH = os.path.join(APP_DIR, 'folders.json')
TRASH_DIR = os.path.join(APP_DIR, '.trash')

CONTENT_CACHE_PATH = os.path.join(APP_DIR, 'content_cache.json')  # --- Content Search Added ---

FILE_CATEGORIES = {
    'Documents': ['.pdf', '.doc', '.docx', '.txt', '.xlsx', '.xls', '.ppt', '.pptx'],
    'Images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'],
    'Videos': ['.mp4', '.avi', '.mov', '.mkv', '.wmv'],
    'Music': ['.mp3', '.wav', '.aac', '.flac', '.m4a']
}

ICON_MAP = {
    'Documents': '📄',
    'Images': '🖼️',
    'Videos': '🎞️',
    'Music': '🎵',
    'Others': '📁'
}


def _shade(hex_color: str, delta: float) -> str:
    """Lighten/darken a #RRGGBB color by delta (-1..+1)."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    def clamp(x): return max(0, min(255, x))
    r = clamp(int(r + (255 - r) * delta) if delta > 0 else int(r * (1 + delta)))
    g = clamp(int(g + (255 - g) * delta) if delta > 0 else int(g * (1 + delta)))
    b = clamp(int(b + (255 - b) * delta) if delta > 0 else int(b * (1 + delta)))
    return f"#{r:02x}{g:02x}{b:02x}"

PALETTE_LIGHT = {
    "bg":        "#f3f5f9",
    "panel":     "#ffffff",
    "altpanel":  "#f4f6f8",
    "muted":     "#6b7280",
    "text":      "#111827",
    "outline":   "#c9ced6",
    "accent":    "#2f6df6",
    "danger":    "#e45555",
}

PALETTE_DARK = {
    "bg":        "#0f1722",
    "panel":     "#0e1a29",
    "altpanel":  "#121e2e",
    "muted":     "#9aa3af",
    "text":      "#e5e7eb",
    "outline":   "#2a3240",
    "accent":    "#3b82f6",
    "danger":    "#ef4444",
}

class SimpleTheme:
    def __init__(self, dark: bool = False):
        self.dark = dark
        self.palette = PALETTE_DARK if dark else PALETTE_LIGHT
    def toggle(self):
        self.dark = not self.dark
        self.palette = PALETTE_DARK if self.dark else PALETTE_LIGHT

def _apply_widget_palette(w: tk.Widget, pal: dict):
    # container background
    for opt in ("bg", "background"):
        try:
            w.configure(**{opt: pal["panel"]})
            break
        except Exception:
            pass

    # text widgets
    if isinstance(w, (tk.Label, tk.Button)):
        try: w.configure(fg=pal["text"])
        except Exception: pass
    if isinstance(w, (tk.Entry,)):
        try:
            w.configure(fg=pal["text"], insertbackground=pal["text"], disabledforeground=pal["muted"])
        except Exception: pass
    if isinstance(w, (tk.Listbox,)):
        try:
            w.configure(
                fg=pal["text"],
                highlightthickness=1,
                highlightbackground=pal["outline"],
                highlightcolor=_shade(pal["outline"], 0.25),
                bd=1, relief="solid",
                selectbackground=_shade(pal["accent"], -0.2),
                selectforeground="#ffffff"
            )
        except Exception: pass


    try:
        for c in w.winfo_children():
            _apply_widget_palette(c, pal)
    except Exception:
        pass

def _style_button(btn: tk.Button, pal: dict, kind: str = "normal"):
    bg = pal["accent"] if kind == "accent" else pal["panel"]
    fg = "#fff" if kind == "accent" else pal["text"]
    hover_bg = _shade(bg, -0.15) if kind == "accent" else _shade(pal["panel"], -0.05)
    outline = pal["outline"]

    btn.configure(
        bg=bg,
        fg=fg,
        bd=0,
        activebackground=hover_bg,
        activeforeground=fg,
        relief="flat",
        highlightthickness=1,
        highlightbackground=outline,
        highlightcolor=_shade(outline, 0.2),
        font=("Segoe UI", 9, "bold"),
        cursor="hand2",
        padx=10, pady=5
    )

    # Simple hover effect
    def on_enter(_): btn.config(bg=hover_bg)
    def on_leave(_): btn.config(bg=bg)
    btn.bind("<Enter>", on_enter)
    btn.bind("<Leave>", on_leave)

def style_entry(e, pal):
    e.configure(
        bg=pal["panel"], fg=pal["text"],
        insertbackground=pal["text"], 
        highlightthickness=1, highlightbackground=pal["outline"],
        highlightcolor=_shade(pal["outline"], 0.25),
        relief="flat", bd=0
    )

def style_optionmenu(om, pal):
    om.configure(
        bg=pal["panel"], fg=pal["text"],
        activebackground=_shade(pal["panel"], -0.05), activeforeground=pal["text"],
        highlightthickness=1, highlightbackground=pal["outline"],
        relief="flat", bd=0
    )

    try:
        menu = om.nametowidget(om.menuname)
        menu.configure(
            bg=pal["panel"], fg=pal["text"],
            activebackground=_shade(pal["panel"], -0.08), activeforeground=pal["text"],
            bd=0, relief="flat", tearoff=False
        )
    except Exception:
        pass

def style_listbox(lb, pal):
    lb.configure(
        bg=pal["panel"], fg=pal["text"],
        selectbackground=_shade(pal["accent"], -0.2), selectforeground="#ffffff",
        highlightthickness=1, highlightbackground=pal["outline"],
        relief="flat", bd=0
    )




def human_size(nbytes: int) -> str:
    units = ['B', 'KB', 'MB', 'GB', 'TB']
    size = float(nbytes)
    i = 0
    while size >= 1024 and i < len(units) - 1:
        size /= 1024
        i += 1
    return f"{size:.2f} {units[i]}"

def ensure_dirs():
    # Make sure trash folder exists
    os.makedirs(TRASH_DIR, exist_ok=True)

    # Automatically create missing JSON files
    required_json = {
        RULES_PATH: [],
        TAGS_PATH: {},
        FOLDERS_PATH: [],
        CONTENT_CACHE_PATH: {}
    }

    for path, default_data in required_json.items():
        if not os.path.exists(path):
            try:
                with open(path, "w", encoding="utf-8") as f:
                    json.dump(default_data, f, indent=2)
            except Exception as e:
                print(f"Error creating {path}: {e}")


def open_in_default(path: str):
    try:
        if not os.path.exists(path):
            messagebox.showwarning("Open File", "File no longer exists.")
            return
        if sys.platform.startswith('darwin'):
            os.system(f"open '{path}'")
        elif os.name == 'nt':
            os.startfile(path) 
        else:
            os.system(f"xdg-open '{path}'")
    except Exception as e:
        messagebox.showerror("Open File", str(e))

def prevent_blank_select(listbox: tk.Listbox):
    def _on_click(e):
        lb = e.widget

        idx = lb.nearest(e.y)
        if idx < 0:
            return "break"
        bbox = lb.bbox(idx) 
        if not bbox:
            return "break"
        x, y, w, h = bbox
        if e.y < y or e.y > y + h:
            lb.selection_clear(0, tk.END)
            return "break"
    listbox.bind("<Button-1>", _on_click, add="+")



# ------------------------------- Data Store -------------------------------
class DataStore:
    def __init__(self):
        self.rules = []
        self.tags = {}
        self.folders = []
        self.undo_stack = []
        self.load_all()

    def load_json(self, path, default):
        if os.path.exists(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                return default
        return default

    def save_json(self, path, data):
        try:
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2)
        except Exception as e:
            messagebox.showerror("Save", f"Failed to save {os.path.basename(path)}: {e}")

    def load_all(self):
        self.rules = self.load_json(RULES_PATH, [])
        self.tags = self.load_json(TAGS_PATH, {})
        self.folders = self.load_json(FOLDERS_PATH, [])

    def save_rules(self):
        self.save_json(RULES_PATH, self.rules)

    def save_tags(self):
        self.save_json(TAGS_PATH, self.tags)

    def save_folders(self):
        self.save_json(FOLDERS_PATH, self.folders)

    def default_folder_path(self):
        for f in self.folders:
            if f.get('default'):
                return f.get('path')
        return None

# --- Content Search Added ---
class ContentSearchEngine:
    SUPPORTED_EXTS = {
        '.txt', '.pdf', '.docx', '.pptx', '.json', '.csv', '.md',
        '.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp' 
    }


    def __init__(self, cache_path=CONTENT_CACHE_PATH):
        self.cache_path = cache_path
        self.cache = self._load_cache()
        self.folder = None
        self.files_meta = []  
        self.texts = []       
        self.vectorizer = None
        self.doc_matrix = None
        self._lock = threading.Lock()
        self._ready = False
        self._building = False
        self._sklearn_ok = (TfidfVectorizer is not None and cosine_similarity is not None)

    def _load_cache(self):
        if os.path.exists(self.cache_path):
            try:
                with open(self.cache_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                return {}
        return {}

    def _save_cache(self):
        try:
            with open(self.cache_path, 'w', encoding='utf-8') as f:
                json.dump(self.cache, f, indent=2)
        except Exception:
            pass 

    def set_folder(self, folder):
        with self._lock:
            self.folder = folder

    def _needs_refresh(self, path, mtime, size):
        c = self.cache.get(path)
        return not c or c.get('mtime') != mtime or c.get('size') != size

    def _extract_text(self, meta):
        path = meta['path']
        ext = meta['ext'].lower()

        try:
            # ---------------- TEXT FILES ----------------
            if ext in ('.txt', '.md'):
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # ---------------- PDF ----------------
            if ext == '.pdf' and fitz is not None:
                text = []
                with fitz.open(path) as doc:
                    for page in doc:
                        text.append(page.get_text("text"))
                return "\n".join(text)

            # ---------------- DOCX ----------------
            if ext == '.docx' and docx is not None:
                d = docx.Document(path)
                parts = [p.text for p in d.paragraphs]

                for table in d.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            parts.append(cell.text)

                return "\n".join(parts)

            # ---------------- PPTX ----------------
            if ext == '.pptx' and Presentation is not None:
                prs = Presentation(path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            parts.append(shape.text)
                return "\n".join(parts)

            # ---------------- JSON ----------------
            if ext == '.json':
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = f.read()
                try:
                    obj = json.loads(data)
                    return json.dumps(obj, ensure_ascii=False, indent=2)
                except Exception:
                    return data

            # ---------------- CSV ----------------
            if ext == '.csv':
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            # --------------------- Image ----------------------
            if ext in ('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.webp'):
                if Image is not None and pytesseract is not None:
                    try:
                        img = Image.open(path)
                        text = pytesseract.image_to_string(img)
                        return text.strip()
                    except Exception as e:
                        print("OCR error:", e)
                        return ""
                else:
                    return ""  # if OCR unavailable  skip image text

        except Exception:
            traceback.print_exc()

        return ""


    def build_index_async(self, files_meta, on_done=None):

        # Build the content index in a background thread to avoid freezing the UI.

        if self._building:
            return
        self._building = True
        self._ready = False

        def worker():
            try:
                texts = []
                for meta in files_meta:
                    path = meta['path']
                    ext = meta['ext'].lower()
                    if ext in self.SUPPORTED_EXTS:
                        if self._needs_refresh(path, meta['mtime'], meta['size']):
                            text = self._extract_text(meta)
                            self.cache[path] = {
                                'mtime': meta['mtime'],
                                'size': meta['size'],
                                'text': text
                            }
                        else:
                            text = self.cache[path].get('text', '')
                        texts.append(text)
                    else:
                        texts.append('')  
                self._save_cache()

                with self._lock:
                    self.files_meta = files_meta[:]
                    self.texts = texts

                    if self._sklearn_ok and any(t.strip() for t in texts):
                        self.vectorizer = TfidfVectorizer(
                            lowercase=True,
                            stop_words='english',
                            ngram_range=(1, 2),
                            max_df=0.95,
                            min_df=1,
                            sublinear_tf=True,
                        )
                        self.doc_matrix = self.vectorizer.fit_transform(texts)
                    else:
                        self.vectorizer = None
                        self.doc_matrix = None

                    self._ready = True
            finally:
                self._building = False
                if callable(on_done):
                    try:
                        on_done()
                    except Exception:
                        pass

        t = threading.Thread(target=worker, daemon=True)
        t.start()

    def query_scores(self, query):
        query = (query or '').strip()
        if not query:
            return {}

        with self._lock:
            if not self._ready or not self.files_meta:
                return {}
            results = {}

            if self.doc_matrix is not None and self.vectorizer is not None:
                try:
                    q_vec = self.vectorizer.transform([query])
                    sims = cosine_similarity(self.doc_matrix, q_vec).ravel()
                    for meta, s in zip(self.files_meta, sims):
                        if s and s > 0:
                            results[meta['path']] = float(s)
                except Exception:
                    pass
            else:
                pattern = re.compile(re.escape(query), re.IGNORECASE)
                for meta, text in zip(self.files_meta, self.texts):
                    if not text:
                        continue
                    matches = pattern.findall(text)
                    if matches:
                        results[meta['path']] = float(len(matches))
            return results

# ------------------------------- Rules Panel -------------------------------
class RulesPanel(tk.Frame):
    def __init__(self, master, store):
        super().__init__(master, bg='#f4f6f8', padx=10, pady=10)
        self.store = store
        self.selected_rules = set()

        tk.Label(self, text='Rules', bg='#f4f6f8', fg='#222',
                 font=('Arial', 12, 'bold')).pack(anchor='w', pady=(0, 8))

        self.rules_list = tk.Listbox(self, activestyle='none', selectmode=tk.SINGLE,
                                     height=18, exportselection=False)
        self.rules_list.pack(fill='both', expand=True)
        self.rules_list.bind('<<ListboxSelect>>', self._on_list_select)
        self.rules_list.bind('<FocusOut>', self._clear_blue_selection)

        #clicking blank space selects last item
        prevent_blank_select(self.rules_list)

        btns = tk.Frame(self, bg='#f4f6f8')
        btns.pack(fill='x', pady=8)

        b_add    = tk.Button(btns, text='Add Rule',    command=self.add_rule,    width=10); b_add.pack(side='left')
        b_edit   = tk.Button(btns, text='Edit Rule',   command=self.edit_rule,   width=10); b_edit.pack(side='left', padx=4)
        b_delete = tk.Button(btns, text='Delete Rule', command=self.delete_rule, width=10); b_delete.pack(side='left')
        self.toggle_btn = tk.Button(btns, text='Select Rule', command=self.toggle_selected_rule, width=12)
        self.toggle_btn.pack(side='left', padx=4)

        self.refresh() 

    def refresh_colors(self, pal: dict):
        for i in range(self.rules_list.size()):
            if i in self.selected_rules:
                self.rules_list.itemconfig(i, bg="#28a745", fg="white")
            else:
                self.rules_list.itemconfig(i, bg=pal["panel"], fg=pal["text"])

    def refresh(self):
        self.rules_list.delete(0, tk.END)
        for i, r in enumerate(self.store.rules):
            name = r.get('name') or self._rule_to_text(r)
            self.rules_list.insert(tk.END, name)
        # use the app's current palette if available
        pal = getattr(getattr(self.master, "master", self.master), "theme", None)
        pal = pal.palette if pal else PALETTE_LIGHT
        self.refresh_colors(pal)
        self._update_toggle_label()


    def get_selected_rules(self):
        return [self.store.rules[i] for i in sorted(self.selected_rules)]

    def current_index(self):
        sel = self.rules_list.curselection()
        return sel[0] if sel else None

    def add_rule(self):
        RuleEditor(self, None, self.store, on_save=self.on_rule_saved)

    def edit_rule(self):
        idx = self.current_index()
        if idx is None:
            messagebox.showinfo('Rules', 'Select a rule to edit.')
            return
        RuleEditor(self, idx, self.store, on_save=self.on_rule_saved)

    def delete_rule(self):
        idx = self.current_index()
        if idx is None:
            messagebox.showinfo('Rules', 'Select a rule to delete.')
            return
        if messagebox.askyesno('Delete Rule', 'Are you sure?'):
            del self.store.rules[idx]
            self.store.save_rules()
            self.selected_rules = {i if i < idx else i - 1 for i in self.selected_rules if i != idx}
            self.refresh()

    def toggle_selected_rule(self):
        idx = self.current_index()
        if idx is None:
            messagebox.showinfo('Rules', 'Select a rule first.')
            return
        if idx in self.selected_rules:
            self.selected_rules.remove(idx)
            self.rules_list.itemconfig(idx, bg='white', fg='black')
        else:
            self.selected_rules.add(idx)
            self.rules_list.itemconfig(idx, bg='#28a745', fg='white')
        self.store.save_rules()
        self._update_toggle_label()

    def _on_list_select(self, _e=None):
        self._update_toggle_label()

    def _clear_blue_selection(self, _e=None):
        self.rules_list.selection_clear(0, tk.END)

    def on_rule_saved(self):
        self.refresh()

    def _update_toggle_label(self):
        idx = self.current_index()
        if idx is not None and idx in self.selected_rules:
            self.toggle_btn.config(text='Unselect Rule')
        else:
            self.toggle_btn.config(text='Select Rule')

    def _rule_to_text(self, r):
        cond = r.get('condition', {})
        action = r.get('action', 'move')
        dest = r.get('destination', '')
        return f"{cond.get('type', '?')} → {action} {('to ' + dest) if dest else ''}"
    
    def apply_palette(self, pal: dict):
        try:
            self.configure(bg=pal["panel"])
        except Exception:
            pass
        _apply_widget_palette(self, pal)
        self.refresh_colors(pal)

        def _walk(w):
            for c in w.winfo_children():
                if isinstance(c, tk.Button):
                    txt = c.cget("text").lower()

                    if any(k in txt for k in ("delete", "remove")):
                        c.config(bg=pal["danger"], fg="white",
                                activebackground=_shade(pal["danger"], -0.15),
                                activeforeground="white", bd=0, relief="flat",
                                highlightthickness=0, cursor="hand2")
                    elif any(k in txt for k in ("run rules", "select folder", "open in files panel", "add rule", "add folder")):

                        _style_button(c, pal, kind="accent")
                    else:
                    
                        _style_button(c, pal, kind="normal")
                _walk(c)
        _walk(self)


    def _style_buttons_recursive(self, widget, pal):
        for c in widget.winfo_children():
            if isinstance(c, tk.Button):
                # highlight "danger" like Delete or Remove
                txt = c.cget("text").lower()
                if "delete" in txt or "remove" in txt:
                    c.config(bg=pal["danger"], fg="white", activebackground=_shade(pal["danger"], -0.15))
                elif "run rules" in txt:
                    _style_button(c, pal, kind="accent")
                else:
                    _style_button(c, pal, kind="normal")
            self._style_buttons_recursive(c, pal)



# ------------------------------- RuleEditor -------------------------------
class RuleEditor(tk.Toplevel):
    def __init__(self, parent, rule_index, store, on_save):
        super().__init__(parent)
        self.title('Rule Editor')
        self.geometry('420x360')
        self.resizable(False, False)
        self.store = store
        self.rule_index = rule_index
        self.on_save = on_save

        self._register_with_main()

        rule = store.rules[rule_index] if rule_index is not None else {
            'name': '',
            'condition': {'type': 'name_contains', 'value': ''},
            'action': 'move',
            'destination': ''
        }

        self.container = tk.Frame(self, padx=10, pady=10)
        self.container.pack(fill='both', expand=True)

        tk.Label(self.container, text='Name').grid(row=0, column=0, sticky='w')
        self.name_var = tk.StringVar(value=rule.get('name', ''))
        self.name_entry = tk.Entry(self.container, textvariable=self.name_var, width=36)
        self.name_entry.grid(row=0, column=1, columnspan=2, sticky='w')

        tk.Label(self.container, text='Condition Type').grid(row=1, column=0, sticky='w', pady=(8, 0))
        self.cond_type = tk.StringVar(value=rule['condition'].get('type', 'name_contains'))
        self.cond_type_menu = tk.OptionMenu(self.container, self.cond_type,
                                            'name_contains', 'extension', 'size_gt_mb', 'size_lt_mb')
        self.cond_type_menu.grid(row=1, column=1, sticky='w', pady=(8, 0))

        tk.Label(self.container, text='Value').grid(row=2, column=0, sticky='w')
        self.cond_value = tk.StringVar(value=str(rule['condition'].get('value', '')))
        self.value_entry = tk.Entry(self.container, textvariable=self.cond_value, width=20)
        self.value_entry.grid(row=2, column=1, sticky='w')
        tk.Label(self.container, text='(text/ext/MB)').grid(row=2, column=2, sticky='w')

        tk.Label(self.container, text='Action').grid(row=3, column=0, sticky='w', pady=(8, 0))
        self.action_var = tk.StringVar(value=rule.get('action', 'move'))
        self.action_menu = tk.OptionMenu(self.container, self.action_var, 'move', 'copy', 'delete')
        self.action_menu.grid(row=3, column=1, sticky='w', pady=(8, 0))

        tk.Label(self.container, text='Destination').grid(row=4, column=0, sticky='w')
        self.dest_var = tk.StringVar(value=rule.get('destination', ''))
        self.dest_entry = tk.Entry(self.container, textvariable=self.dest_var, width=28)
        self.dest_entry.grid(row=4, column=1, sticky='w')
        self.browse_btn = tk.Button(self.container, text='Browse', command=self.browse_dest)
        self.browse_btn.grid(row=4, column=2, sticky='w')

        btns = tk.Frame(self.container)
        btns.grid(row=5, column=0, columnspan=3, pady=12, sticky='e')
        self.btn_save = tk.Button(btns, text='Save & Close', command=self.save)
        self.btn_save.pack(side='right')
        self.btn_cancel = tk.Button(btns, text='Cancel', command=self.destroy)
        self.btn_cancel.pack(side='right', padx=6)

        try:
            pal = self._main().theme.palette
        except Exception:
            pal = PALETTE_LIGHT
        self.apply_palette(pal)

    # --- theme ---
    def _main(self):
        return self.master.winfo_toplevel()

    def _register_with_main(self):
        try:
            root = self._main()
            if not hasattr(root, "_editors"):
                root._editors = set()
            root._editors.add(self)
        except Exception:
            pass

    def destroy(self):
        try:
            root = self._main()
            if hasattr(root, "_editors") and self in root._editors:
                root._editors.remove(self)
        except Exception:
            pass
        super().destroy()

    def apply_palette(self, pal: dict):
        self.configure(bg=pal["panel"])
        self.container.configure(bg=pal["panel"])

        # labels
        for w in self.container.winfo_children():
            if isinstance(w, tk.Label):
                w.configure(bg=pal["panel"], fg=pal["text"])

        # inputs/menus
        style_entry(self.name_entry, pal)
        style_entry(self.value_entry, pal)
        style_entry(self.dest_entry, pal)
        style_optionmenu(self.cond_type_menu, pal)
        style_optionmenu(self.action_menu, pal)

        # buttons
        _style_button(self.btn_save, pal, kind="accent")
        _style_button(self.btn_cancel, pal, kind="normal")
        _style_button(self.browse_btn, pal, kind="normal")


    def browse_dest(self):
        path = filedialog.askdirectory(title='Choose destination folder')
        if path:
            self.dest_var.set(path)

    def save(self):
        name = self.name_var.get().strip()
        cond_type = self.cond_type.get()
        cond_val = self.cond_value.get().strip()
        action = self.action_var.get()
        dest = self.dest_var.get().strip()

        if action in ('move', 'copy') and not dest:
            default = self.store.default_folder_path()
            if default:
                dest = default
            else:
                messagebox.showerror('Rule', 'Destination folder required.')
                return

        rule = {
            'name': name or f"{cond_type} {cond_val}",
            'condition': {'type': cond_type, 'value': cond_val},
            'action': action,
            'destination': dest
        }
        if self.rule_index is None:
            self.store.rules.append(rule)
        else:
            self.store.rules[self.rule_index] = rule
        self.store.save_rules()
        if self.on_save:
            self.on_save()
        self.destroy()

# ------------------------------- File Explorer Panel -------------------------------
class FileExplorerPanel(tk.Frame):
    def __init__(self, master, store):
        super().__init__(master, bg='#ffffff', padx=10, pady=10)
        self.store = store
        self.current_folder = None
        self.files = []
        self.filtered_indices = []
        self.last_batch = []
        self._rules_provider = None
        self.preview_panel = None

        # --- Content Search Added ---
        self.search_engine = ContentSearchEngine()
        self._index_status_var = tk.StringVar(value='') 


        tk.Label(self, text='Files', bg='#ffffff', fg='#222', font=('Arial', 12, 'bold')).pack(anchor='w', pady=(0, 8))

        self.toolbar = tk.Frame(self, bg='#eef1f5', padx=8, pady=6)
        self.toolbar.pack(fill='x', pady=(0, 8))
        tk.Button(self.toolbar, text='Select Folder', command=self.select_folder).pack(side='left')


        tk.Label(self.toolbar, text=' Filter: ', bg='#eef1f5').pack(side='left')
        self.filter_var = tk.StringVar(value='All')
        self.filter_menu = tk.OptionMenu(self.toolbar, self.filter_var, 'All', 'Documents', 'Images', 'Videos', 'Music', 'Others',
                                        command=lambda _: self.refresh_list())
        self.filter_menu.pack(side='left')

        tk.Label(self.toolbar, text=' Search: ', bg='#eef1f5').pack(side='left')
        self.search_var = tk.StringVar()
        self.search_entry = tk.Entry(self.toolbar, textvariable=self.search_var, width=18)
        self.search_entry.pack(side='left')
        self.search_entry.bind('<KeyRelease>', lambda _e: self.refresh_list())

        tk.Label(self.toolbar, text=' Tag: ', bg='#eef1f5').pack(side='left')
        self.tag_filter_var = tk.StringVar()
        self.tag_filter_entry = tk.Entry(self.toolbar, textvariable=self.tag_filter_var, width=14)
        self.tag_filter_entry.pack(side='left')
        self.tag_filter_entry.bind('<KeyRelease>', lambda _e: self.refresh_list())

        tk.Label(self.toolbar, text=' Sort by: ', bg='#eef1f5').pack(side='left')
        self.sort_var = tk.StringVar(value='Name')
        self.sort_menu = tk.OptionMenu(self.toolbar, self.sort_var, 'Name', 'Type', 'Date Modified', 'Size',
                                    command=lambda _: self.refresh_list())
        self.sort_menu.pack(side='left')


        main_area = tk.Frame(self, bg='#ffffff')
        main_area.pack(fill='both', expand=True)
        self.file_list = tk.Listbox(main_area, activestyle='none', selectmode=tk.EXTENDED)
        self.file_list.pack(side='left', fill='both', expand=True)
        sb = tk.Scrollbar(main_area, orient='vertical', command=self.file_list.yview)
        sb.pack(side='right', fill='y')
        self.file_list.config(yscrollcommand=sb.set)
        self.file_list.bind('<Double-Button-1>', self.open_selected)
        self.file_list.bind('<<ListboxSelect>>', self._on_select_change)


 
        prevent_blank_select(self.file_list)

        tagbar = tk.Frame(self, bg='#ffffff')
        tagbar.pack(fill='x', pady=6)
        tk.Label(tagbar, text='Tag: ', bg='#ffffff').pack(side='left')
        self.tag_entry_var = tk.StringVar()
        self.tag_entry = tk.Entry(tagbar, textvariable=self.tag_entry_var, width=20)
        self.tag_entry.pack(side='left')

        tk.Button(tagbar, text='Add Tag to Selected', command=self.add_tag_selected).pack(side='left', padx=4)
        tk.Button(tagbar, text='Remove Tag from Selected', command=self.remove_tag_selected).pack(side='left')

        actions = tk.Frame(self, bg='#ffffff')
        actions.pack(fill='x', pady=(8, 0))
        tk.Button(actions, text='Run Rules', command=self.run_rules_here).pack(side='left')
        tk.Button(actions, text='Undo Last Action', command=self.undo_last).pack(side='left', padx=6)

        self.info_var = tk.StringVar(value='No folder selected.')
        tk.Label(self, textvariable=self.info_var, bg='#ffffff', fg='#555').pack(anchor='w', pady=(6, 0))

  
        tk.Label(self, textvariable=self._index_status_var, bg='#ffffff', fg='#888', font=('Arial', 9)).pack(anchor='w')


    def set_rules_provider(self, provider_callable):
        self._rules_provider = provider_callable

    # ---------- Folder & file loading ----------
    def select_folder(self):
        folder = filedialog.askdirectory(title='Select a folder')
        if folder:
            self.load_folder(folder)

    def load_folder(self, folder):
        self.current_folder = folder
        self.files = []
        try:
            with os.scandir(folder) as it:
                for entry in it:
                    if entry.is_file():
                        stat = entry.stat()
                        self.files.append({
                            'name': entry.name,
                            'path': entry.path,
                            'ext': Path(entry.name).suffix.lower(),
                            'size': stat.st_size,
                            'mtime': stat.st_mtime
                        })
        except Exception as e:
            messagebox.showerror('Folder', str(e))
            return
        self.refresh_list()
        self.info_var.set(f"{folder} • {len(self.files)} files")

        # --- Content Search Added ---
        # Build index in the background.
        self._index_status_var.set('Indexing contents for semantic search…')
        self.search_engine.set_folder(folder)
        self.search_engine.build_index_async(self.files, on_done=self._on_index_ready)

 
    def _on_index_ready(self):
        self._index_status_var.set('Content index ready.')
        # Re-run refresh to include content scores
        try:
            self.refresh_list()
        except Exception:
            pass


    def refresh_list(self):
        self.file_list.delete(0, tk.END)
        if not self.files:
            self.filtered_indices = []
            return

        rows = list(self.files)
        fcat = self.filter_var.get()
        if fcat != 'All':
            if fcat == 'Others':
                rows = [r for r in rows if not any(r['ext'] in exts for exts in FILE_CATEGORIES.values())]
            else:
                rows = [r for r in rows if r['ext'] in FILE_CATEGORIES.get(fcat, [])]

        q = self.search_var.get().strip().lower()

        # --- Content Search Added ---
        # compute content-based scores once (dict path - score)
        content_scores = {}
        if q:
            content_scores = self.search_engine.query_scores(q)
        # --- Content Search Added ---

        if q:
            # filename match OR content match
            def base_name_score(r):
                name = r['name'].lower()
                if q in name:
                    #filename relevance prefix, then substring
                    if name.startswith(q):
                        return 1.0
                    return 0.6
                return 0.0

            # merge: filename score + content score 
            scored = []
            for r in rows:
                s_name = base_name_score(r)
                s_content = content_scores.get(r['path'], 0.0)
                score = s_name + (1.2 * s_content)
                if s_name > 0 or s_content > 0:
                    scored.append((score, r))

            # If none matched fall back to filename filter
            if scored:
                scored.sort(key=lambda x: x[0], reverse=True)
                rows = [r for _s, r in scored]
            else:
                # fallback to original filename filter behavior
                rows = [r for r in rows if q in r['name'].lower()]
        # else no query proceed as original

        tq = self.tag_filter_var.get().strip().lower()
        if tq:
            rows = [
                r for r in rows
                if any(tq in t.lower() for t in self.store.tags.get(r['path'], []))
            ]

        # Keep user sort when no query.
        if not q:
            key = self.sort_var.get()
            if key == 'Name':
                rows.sort(key=lambda r: r['name'].lower())
            elif key == 'Type':
                rows.sort(key=lambda r: r['ext'])
            elif key == 'Date Modified':
                rows.sort(key=lambda r: r['mtime'], reverse=True)
            elif key == 'Size':
                rows.sort(key=lambda r: r['size'], reverse=True)

        self.filtered_indices = [self.files.index(r) for r in rows]

        for r in rows:
            tags = ','.join(self.store.tags.get(r['path'], []))
            dt = datetime.fromtimestamp(r['mtime']).strftime('%Y-%m-%d %H:%M')
            ic = ICON_MAP['Others']
            for cat, exts in FILE_CATEGORIES.items():
                if r['ext'] in exts:
                    ic = ICON_MAP.get(cat, ICON_MAP['Others'])
                    break
            display = f"{ic}  {r['name']}  | {r['ext'] or '-'} | {human_size(r['size'])} | {dt} | tags:[{tags}]"
            self.file_list.insert(tk.END, display)

    # ---------- Selection & tagging ----------
    def selected_paths(self):
        paths = []
        for idx in self.file_list.curselection():
            if 0 <= idx < len(self.filtered_indices):
                original = self.filtered_indices[idx]
                paths.append(self.files[original]['path'])
        return paths

    # Case-insensitive Add Tag
    def add_tag_selected(self):
        tag = self.tag_entry_var.get().strip()
        if not tag:
            return
        tag_lower = tag.lower()
        any_added = False
        for p in self.selected_paths():
            tags = self.store.tags.get(p, [])
            # normalize and check lowercase duplicates
            tags_lower = [t.lower() for t in tags]
            if tag_lower not in tags_lower:
                tags.append(tag)
                self.store.tags[p] = tags
                any_added = True
        if any_added:
            self.store.save_tags()
            self.refresh_list()
            self.tag_entry_var.set('')

    # Case-insensitive Remove Tag
    def remove_tag_selected(self):
        tag = self.tag_entry_var.get().strip()
        if not tag:
            return
        tag_lower = tag.lower()
        any_removed = False
        for p in self.selected_paths():
            tags = self.store.tags.get(p, [])
            new_tags = [t for t in tags if t.lower() != tag_lower]
            if len(new_tags) != len(tags):
                self.store.tags[p] = new_tags
                any_removed = True
        if any_removed:
            self.store.save_tags()
            self.refresh_list()
            self.tag_entry_var.set('')

    # ---------- File actions ----------
    def open_selected(self, _e=None):
        paths = self.selected_paths()
        if paths:
            open_in_default(paths[0])

    def run_rules_here(self):
        if not self.current_folder:
            default_path = self.store.default_folder_path()
            if default_path:
                self.load_folder(default_path)
            else:
                messagebox.showinfo('Rules', 'Select a folder first.')
                return

        selected_rules = []
        if self._rules_provider is not None:
            try:
                selected_rules = self._rules_provider()
            except Exception:
                selected_rules = []

        rules_to_apply = selected_rules if selected_rules else None
        self.apply_rules_to_folder(self.current_folder, rules=rules_to_apply)

    def apply_rules_to_folder(self, folder, rules=None):
        ensure_dirs()
        if rules is None:
            rules = self.store.rules
        if not rules:
            messagebox.showinfo('Rules', 'No rules to apply.')
            return
        batch = []
        file_entries = [f for f in self.files]
        for r in file_entries:
            path = r['path']
            for rule in rules:
                if self._match_rule(r, rule):
                    action = rule.get('action', 'move')
                    dest = rule.get('destination', '')
                    if action == 'move':
                        os.makedirs(dest, exist_ok=True)
                        new_path = self._unique_path(os.path.join(dest, os.path.basename(path)))
                        shutil.move(path, new_path)
                        batch.append({'type': 'move', 'src': path, 'dst': new_path})
                    elif action == 'copy':
                        os.makedirs(dest, exist_ok=True)
                        new_path = self._unique_path(os.path.join(dest, os.path.basename(path)))
                        shutil.copy2(path, new_path)
                        batch.append({'type': 'copy', 'src': path, 'dst': new_path})
                    elif action == 'delete':
                        os.makedirs(TRASH_DIR, exist_ok=True)
                        new_path = self._unique_path(os.path.join(TRASH_DIR, os.path.basename(path)))
                        shutil.move(path, new_path)
                        batch.append({'type': 'trash', 'src': path, 'dst': new_path})
                    break
        if batch:
            self.store.undo_stack.append(batch)
        self.load_folder(folder)
        messagebox.showinfo('Rules', f'Applied rules. Operations: {len(batch)}')

    def undo_last(self):
        if not self.store.undo_stack:
            messagebox.showinfo('Undo', 'Nothing to undo.')
            return
        batch = self.store.undo_stack.pop()
        for op in reversed(batch):
            typ = op['type']
            src = op['src']
            dst = op['dst']
            if typ == 'move':
                os.makedirs(os.path.dirname(src), exist_ok=True)
                shutil.move(dst, self._unique_path(src))
            elif typ == 'copy':
                if os.path.exists(dst):
                    os.remove(dst)
            elif typ == 'trash':
                os.makedirs(os.path.dirname(src), exist_ok=True)
                shutil.move(dst, self._unique_path(src))
        if self.current_folder:
            self.load_folder(self.current_folder)
        messagebox.showinfo('Undo', 'Undo complete.')

    def _unique_path(self, path):
        if not os.path.exists(path):
            return path
        base, ext = os.path.splitext(path)
        i = 1
        while True:
            cand = f"{base} ({i}){ext}"
            if not os.path.exists(cand):
                return cand
            i += 1

    def _match_rule(self, file_row, rule):
        c = rule.get('condition', {})
        ctype = c.get('type')
        val = c.get('value', '')
        if ctype == 'name_contains':
            return str(val).lower() in file_row['name'].lower()
        if ctype == 'extension':
            v = val.lower().strip()
            if v and not v.startswith('.'):
                v = '.' + v
            return file_row['ext'] == v
        if ctype == 'size_gt_mb':
            try:
                return file_row['size'] > float(val) * 1024 * 1024
            except Exception:
                return False
        if ctype == 'size_lt_mb':
            try:
                return file_row['size'] < float(val) * 1024 * 1024
            except Exception:
                return False
        return False
    
    def apply_palette(self, pal: dict):
        # panel bg
        try:
            self.configure(bg=pal["panel"])
        except Exception:
            pass

        _apply_widget_palette(self, pal)

        # toolbar bg reacts to theme
        try:
            self.toolbar.configure(bg=pal["altpanel"])
        except Exception:
            pass
        # fix the small text labels on the toolbar to match
        for child in self.toolbar.winfo_children():
            if isinstance(child, tk.Label):
                child.configure(bg=pal["altpanel"], fg=pal["text"])

        # theming so dark mode stays readable
        try:
            style_entry(self.search_entry, pal)
            style_entry(self.tag_filter_entry, pal)
            style_entry(self.tag_entry, pal)
            style_optionmenu(self.filter_menu, pal)
            style_optionmenu(self.sort_menu, pal)
            style_listbox(self.file_list, pal)
        except Exception:
            pass

        def _walk(w):
            for c in w.winfo_children():
                if isinstance(c, tk.Button):
                    txt = c.cget("text").lower()
                    if any(k in txt for k in ("delete", "remove")):
                        c.config(bg=pal["danger"], fg="white",
                                activebackground=_shade(pal["danger"], -0.15),
                                activeforeground="white", bd=0, relief="flat",
                                highlightthickness=0, cursor="hand2")
                    elif any(k in txt for k in ("run rules", "select folder", "open in files panel", "add rule", "add folder")):
                        _style_button(c, pal, kind="accent")
                    else:
                        _style_button(c, pal, kind="normal")
                _walk(c)
        _walk(self)


    def _style_buttons_recursive(self, widget, pal):
        for c in widget.winfo_children():
            if isinstance(c, tk.Button):
                # highlight "danger" like Delete or Remove
                txt = c.cget("text").lower()
                if "delete" in txt or "remove" in txt:
                    c.config(bg=pal["danger"], fg="white", activebackground=_shade(pal["danger"], -0.15))
                elif "run rules" in txt:
                    _style_button(c, pal, kind="accent")
                else:
                    _style_button(c, pal, kind="normal")
            self._style_buttons_recursive(c, pal)


    def set_preview_panel(self, panel):
        self.preview_panel = panel

    
    def _on_select_change(self, _e=None):
        if not self.preview_panel:
            return
        paths = self.selected_paths()
        if paths:
            self.preview_panel.update_preview(paths[0])


# ------------------------------- Preview Panel -------------------------------
class PreviewPanel(tk.Frame):
    def __init__(self, master):
        super().__init__(master, bg="#f7f7f7", padx=10, pady=10)

        # for image resizing keeping references
        self._photo = None         
        self._photos = []          
        self._orig_image = None    

        self.title_var = tk.StringVar(value="Preview")
        self.title_label = tk.Label(
            self, textvariable=self.title_var,
            font=("Arial", 12, "bold"), bg="#f7f7f7", fg="#111827", anchor="w"
        )
        self.title_label.pack(anchor="w", pady=(0, 4), fill="x")

        # outer frame
        self.content_frame = tk.Frame(self, bg="#ffffff", bd=1, relief="solid")
        self.content_frame.pack(fill="both", expand=True)

        # re-draw image on resize so it always fits nicely
        self.content_frame.bind("<Configure>", self._on_content_resize)

        # canvas + scrollbars
        self.canvas = tk.Canvas(
            self.content_frame,
            bg="#ffffff", highlightthickness=0, bd=0
        )
        self.vscroll = tk.Scrollbar(
            self.content_frame, orient="vertical",
            command=self.canvas.yview
        )
        self.hscroll = tk.Scrollbar(
            self.content_frame, orient="horizontal",
            command=self.canvas.xview
        )
        self.canvas.configure(
            yscrollcommand=self.vscroll.set,
            xscrollcommand=self.hscroll.set
        )

        # grid so scrollbars never disappear
        self.content_frame.grid_rowconfigure(0, weight=1)
        self.content_frame.grid_rowconfigure(1, weight=0)
        self.content_frame.grid_columnconfigure(0, weight=1)
        self.content_frame.grid_columnconfigure(1, weight=0)

        self.canvas.grid(row=0, column=0, sticky="nsew")
        self.vscroll.grid(row=0, column=1, sticky="ns")
        self.hscroll.grid(row=1, column=0, sticky="ew")

        # inner frame that actually holds the content
        self.inner = tk.Frame(self.canvas, bg="#ffffff")
        self.inner_window = self.canvas.create_window(
            (0, 0), window=self.inner, anchor="nw"
        )

        # track size changes so scrolling works properly
        self.inner.bind("<Configure>", self._on_inner_configure)
        self.canvas.bind("<Configure>", self._on_canvas_configure)

        # Initial message
        self.info_label = tk.Label(
            self.inner, text="Select a file to preview.",
            bg="#ffffff", fg="#555555", justify="left", anchor="nw"
        )
        self.info_label.pack(anchor="nw", padx=8, pady=8)

    # ---------- Scroll helpers ----------
    def _on_inner_configure(self, event):
        self.canvas.configure(scrollregion=self.canvas.bbox("all"))

    def _on_canvas_configure(self, event):
        self.canvas.configure(scrollregion=self.canvas.bbox("all"))

    # ---------- Theming ----------
    def apply_palette(self, pal: dict):
        try:
            self.configure(bg=pal["altpanel"])
        except Exception:
            pass

        self.title_label.configure(bg=pal["altpanel"], fg=pal["text"])
        self.content_frame.configure(bg=pal["panel"])
        self.canvas.configure(bg=pal["panel"])
        self.inner.configure(bg=pal["panel"])

        if isinstance(self.info_label, tk.Label):
            self.info_label.configure(bg=pal["panel"], fg=pal["muted"])

        try:
            self.vscroll.configure(troughcolor=pal["panel"])
            self.hscroll.configure(troughcolor=pal["panel"])
        except Exception:
            pass

    # ---------- Helpers ----------
    def _clear_content(self):
        for w in self.inner.winfo_children():
            w.destroy()
        self._photos = []
        self._photo = None

    def _short_name(self, name: str, max_len: int = 40) -> str:
        if len(name) <= max_len:
            return name
        return name[:20] + "..." + name[-15:]

    def show_message(self, msg: str):
        self._clear_content()
        self.title_var.set("Preview")
        self.info_label = tk.Label(
            self.inner, text=msg,
            bg=self.inner["bg"], fg="#555555",
            justify="left", anchor="nw", wraplength=600
        )
        self.info_label.pack(anchor="nw", padx=8, pady=8)

    # ---------- Resize logic ----------
    def _resize_image_to_fit(self, img):
        self.canvas.update_idletasks()

        max_w = max(self.canvas.winfo_width(), 50)
        max_h = max(self.canvas.winfo_height(), 50)

        img_w, img_h = img.size

        # Scale to FIT WIDTH first
        scale_w = max_w / img_w
        scale_h = max_h / img_h

        scale = min(scale_w, scale_h)

        new_size = (int(img_w * scale), int(img_h * scale))
        return img.resize(new_size, Image.LANCZOS)


    def _render_single_image(self, path: str):
        """Render the current _orig_image nicely centered & scaled."""
        if self._orig_image is None:
            return

        self._clear_content()
        base = os.path.basename(path) if path else ""
        if base:
            self.title_var.set(f"Preview - {self._short_name(base)}")
        else:
            pass

        resized = self._resize_image_to_fit(self._orig_image)
        photo = ImageTk.PhotoImage(resized)
        self._photo = photo
        self._photos = [photo]

        lbl = tk.Label(self.inner, image=photo, bg=self.inner["bg"])
        lbl.pack(anchor="center", pady=8)

    def _on_content_resize(self, event=None):
        # re-render it to fit the new size.
        if self._orig_image is not None:
            # pass empty path so don't change the title text
            self._render_single_image("")

    # ---------- Text preview ----------
    def show_text_preview(self, path: str, content: str):
        self._orig_image = None  
        self._clear_content()
        base = os.path.basename(path)
        self.title_var.set(f"Preview - {self._short_name(base)}")

        txt = tk.Text(
            self.inner, wrap="word", state="normal",
            bg=self.inner["bg"], fg="#111827",
            relief="flat", bd=0
        )
        txt.insert("1.0", content)
        txt.config(state="disabled")
        txt.pack(fill="both", expand=True, padx=8, pady=8)

    # ---------- PDF pages----------
    def show_pages_preview(self, path: str, pil_images):
        # used for PDFs
        self._orig_image = None  
        self._clear_content()
        base = os.path.basename(path)
        self.title_var.set(f"Preview - {self._short_name(base)}")

        for i, img in enumerate(pil_images, start=1):
            resized = self._resize_image_to_fit(img)
            photo = ImageTk.PhotoImage(resized)
            self._photos.append(photo)

            lbl = tk.Label(self.inner, image=photo, bg=self.inner["bg"])
            lbl.pack(anchor="center", pady=(8 if i == 1 else 4, 8))

    # ---------- Main entry ----------
    def update_preview(self, path: str):
        if not path or not os.path.exists(path):
            self.show_message("File not found.")
            return

        ext = Path(path).suffix.lower()

        # Normal images 
        if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp") and Image is not None and ImageTk is not None:
            try:
                self._orig_image = Image.open(path)
                self._render_single_image(path)
            except Exception as e:
                self.show_message(f"Image preview error:\n{e}")
            return

        # PDF
        if ext == ".pdf" and fitz is not None and Image is not None and ImageTk is not None:
            try:
                pil_pages = []
                max_pages = 8
                zoom_matrix = fitz.Matrix(2.5, 2.5)

                with fitz.open(path) as doc:
                    page_count = min(doc.page_count, max_pages)
                    if page_count == 0:
                        self.show_message("[Empty PDF]")
                        return

                    for i in range(page_count):
                        page = doc.load_page(i)
                        pix = page.get_pixmap(matrix=zoom_matrix)

                        mode = "RGBA" if pix.alpha else "RGB"
                        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                        pil_pages.append(img)

                self.show_pages_preview(path, pil_pages)

            except Exception as e:
                self.show_message(f"PDF preview error:\n{e}")
            return



        # PPTX text from all slides
        if ext == ".pptx" and Presentation is not None:
            try:
                prs = Presentation(path)
                parts = []
                for idx, slide in enumerate(prs.slides, start=1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            if shape.text.strip():
                                slide_texts.append(shape.text.strip())
                    if slide_texts:
                        parts.append(f"--- Slide {idx} ---\n" + "\n\n".join(slide_texts))
                text = "\n\n".join(parts).strip() or "[No text found in slides]"
                self.show_text_preview(path, text)
            except Exception as e:
                self.show_message(f"PPTX preview error:\n{e}")
            return

        # DOCX full document text
        if ext == ".docx" and docx is not None:
            try:
                d = docx.Document(path)
                parts = [p.text for p in d.paragraphs if p.text.strip()]
                text = "\n\n".join(parts) or "[No text found in document]"
                self.show_text_preview(path, text)
            except Exception as e:
                self.show_message(f"DOCX preview error:\n{e}")
            return

        # Text files
        if ext in (".txt", ".md", ".py", ".json", ".csv", ".log"):
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    data = f.read()
                self.show_text_preview(path, data)
            except Exception as e:
                self.show_message(f"Text preview error:\n{e}")
            return

        # Fallback
        self.show_message(f"No preview available for:\n{os.path.basename(path)}")



# ------------------------------- Folder Panel -------------------------------
class FolderPanel(tk.Frame):
    def __init__(self, master, store, on_choose_for_files):
        super().__init__(master, bg='#f7f7f7', padx=10, pady=10)
        self.store = store
        self.on_choose_for_files = on_choose_for_files

        tk.Label(self, text='Folders', bg='#f7f7f7', fg='#222', font=('Arial', 12, 'bold')).pack(anchor='w', pady=(0,8))

        self.listbox = tk.Listbox(self, activestyle='none', selectmode=tk.SINGLE, height=16)
        self.listbox.pack(fill='both', expand=True)
        self.listbox.bind('<<ListboxSelect>>', lambda e: self.refresh_details())
        
        # bug fix
        prevent_blank_select(self.listbox)

        btns = tk.Frame(self, bg='#f7f7f7')
        btns.pack(fill='x', pady=8)
        tk.Button(btns, text='Add Folder', command=self.add_folder).pack(side='left')
        tk.Button(btns, text='Edit Name', command=self.edit_folder).pack(side='left', padx=4)
        tk.Button(btns, text='Remove', command=self.remove_folder).pack(side='left')

        btns2 = tk.Frame(self, bg='#f7f7f7')
        btns2.pack(fill='x', pady=(0,8))
        tk.Button(btns2, text='Choose Path', command=self.choose_path).pack(side='left')
        self.default_btn = tk.Button(btns2, text='Set as Default', command=self.toggle_default)
        self.default_btn.pack(side='left', padx=4)

    
        self.details = tk.Frame(self, bg='#eef1f5', padx=8, pady=6, bd=1, relief='solid', highlightthickness=1, highlightbackground='#cbd3dc')
        self.details.pack(fill='x')
        self.path_var = tk.StringVar()
        self.count_var = tk.StringVar(value='Files: -')
        self.size_var = tk.StringVar(value='Total size: -')
        tk.Label(self.details, textvariable=self.path_var, bg='#eef1f5').pack(anchor='w')
        tk.Label(self.details, textvariable=self.count_var, bg='#eef1f5').pack(anchor='w')
        tk.Label(self.details, textvariable=self.size_var, bg='#eef1f5').pack(anchor='w')

        tk.Button(self, text='Open in Files Panel', command=self.open_in_files).pack(anchor='e', pady=8)

        self.refresh()

    def refresh(self):
        self.listbox.delete(0, tk.END)
        for f in self.store.folders:
            name = f.get('name') or f.get('path') or '(unnamed)'
            if f.get('default'):
                name += '  [default]'
            self.listbox.insert(tk.END, name)
        self.refresh_details()

    def idx(self):
        sel = self.listbox.curselection()
        return sel[0] if sel else None

    def add_folder(self):
        name = self.simple_input('Folder Name', 'Enter a short name:')
        if not name:
            return
        path = filedialog.askdirectory(title='Choose folder path')
        if not path:
            return
        self.store.folders.append({'name': name, 'path': path, 'default': False})
        self.store.save_folders()
        self.refresh()

    def edit_folder(self):
        i = self.idx()
        if i is None:
            messagebox.showinfo('Folders','Select an item.')
            return
        name = self.simple_input('Edit Name', 'New name:', self.store.folders[i].get('name',''))
        if name:
            self.store.folders[i]['name'] = name
            self.store.save_folders()
            self.refresh()

    def remove_folder(self):
        i = self.idx()
        if i is None:
            messagebox.showinfo('Folders','Select an item.')
            return
        if messagebox.askyesno('Remove','Are you sure?'):
            del self.store.folders[i]
            self.store.save_folders()
            self.refresh()

    def choose_path(self):
        i = self.idx()
        if i is None:
            messagebox.showinfo('Folders','Select an item.')
            return
        path = filedialog.askdirectory(title='Choose folder path')
        if path:
            self.store.folders[i]['path'] = path
            self.store.save_folders()
            self.refresh()

    def toggle_default(self):
        i = self.idx()
        if i is None:
            messagebox.showinfo('Folders','Select an item.')
            return
        folder = self.store.folders[i]
        if folder.get('default'):
            folder['default'] = False
        else:
            for f in self.store.folders:
                f['default'] = False
            folder['default'] = True
        self.store.save_folders()
        self.refresh()

    def refresh_details(self):
        i = self.idx()
        if i is None:
            self.path_var.set('')
            self.count_var.set('Files: -')
            self.size_var.set('Total size: -')
            self.default_btn.config(text='Set as Default')
            return
        path = self.store.folders[i].get('path')
        self.path_var.set(f"Path: {path}")
        count, total = self._folder_stats(path)
        self.count_var.set(f"Files: {count}")
        self.size_var.set(f"Total size: {human_size(total)}")
        if self.store.folders[i].get('default'):
            self.default_btn.config(text='Remove as Default')
        else:
            self.default_btn.config(text='Set as Default')

    def _folder_stats(self, folder):
        count = 0
        total = 0
        try:
            with os.scandir(folder) as it:
                for e in it:
                    if e.is_file():
                        count += 1
                        total += e.stat().st_size
        except Exception:
            pass
        return count, total

    def open_in_files(self):
        i = self.idx()
        if i is None:
            return
        path = self.store.folders[i].get('path')
        if path:
            self.on_choose_for_files(path)

    def simple_input(self, title, prompt, initial=''):
        dlg = tk.Toplevel(self)
        dlg.title(title)
        dlg.geometry('320x120')
        dlg.resizable(False, False)
        tk.Label(dlg, text=prompt).pack(pady=(12,6))
        var = tk.StringVar(value=initial)
        e = tk.Entry(dlg, textvariable=var)
        e.pack(fill='x', padx=12)
        out = {'value': None}
        def ok():
            out['value'] = var.get().strip()
            dlg.destroy()
        btns = tk.Frame(dlg)
        btns.pack(pady=8)
        tk.Button(btns, text='OK', command=ok).pack(side='left')
        tk.Button(btns, text='Cancel', command=dlg.destroy).pack(side='left', padx=6)
        e.focus_set()
        dlg.transient(self)
        dlg.grab_set()
        self.wait_window(dlg)
        return out['value']
    
    def apply_palette(self, pal: dict):
        try:
            self.configure(bg=pal["altpanel"])
        except Exception:
            pass

        _apply_widget_palette(self, pal)

        # listbox (folders)
        try:
            style_listbox(self.listbox, pal)
        except Exception:
            pass

        try:
            self.details.configure(
                bg=pal["panel"],
                highlightthickness=1, highlightbackground=pal["outline"],
                bd=0, relief="flat"
            )
            for child in self.details.winfo_children():
                if isinstance(child, tk.Label):
                    child.configure(bg=pal["panel"], fg=pal["text"])
        except Exception:
            pass

        def _walk(w):
            for c in w.winfo_children():
                if isinstance(c, tk.Button):
                    txt = c.cget("text").lower()
                    if any(k in txt for k in ("delete", "remove")):
                        c.config(bg=pal["danger"], fg="white",
                                activebackground=_shade(pal["danger"], -0.15),
                                activeforeground="white", bd=0, relief="flat",
                                highlightthickness=0, cursor="hand2")
                    elif any(k in txt for k in ("run rules", "select folder", "open in files panel", "add rule", "add folder")):
                        _style_button(c, pal, kind="accent")
                    else:
                        _style_button(c, pal, kind="normal")
                _walk(c)
        _walk(self)


    def _style_buttons_recursive(self, widget, pal):
        for c in widget.winfo_children():
            if isinstance(c, tk.Button):
                txt = c.cget("text").lower()
                if "delete" in txt or "remove" in txt:
                    c.config(bg=pal["danger"], fg="white", activebackground=_shade(pal["danger"], -0.15))
                elif "run rules" in txt:
                    _style_button(c, pal, kind="accent")
                else:
                    _style_button(c, pal, kind="normal")
            self._style_buttons_recursive(c, pal)



# ------------------------------- Main App -------------------------------
class MainApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title('File Manager')
        self.geometry('1500x900')
        self.minsize(1450, 780)
        ensure_dirs()

        self.theme = SimpleTheme(dark=False)  # start in light mode
        pal = self.theme.palette

        # Top bar darker than bg
        self.topbar = tk.Frame(self, bg=_shade(pal["bg"], -0.25), height=56, highlightthickness=0, bd=0)
        self.topbar.pack(fill='x')

        # title group
        title_wrap = tk.Frame(self.topbar, bg=_shade(pal["bg"], -0.25), highlightthickness=0, bd=0)
        title_wrap.pack(side='left', padx=12, pady=8)

        self.title_label = tk.Label(title_wrap, text="🗂️ File Manager",
                                    bg=title_wrap['bg'], fg=pal["text"], font=("Segoe UI", 12, "bold"))
        self.title_label.pack(side='left')

        self.subtitle_label = tk.Label(title_wrap, text="  organize • tag • search",
                                    bg=title_wrap['bg'], fg=pal["muted"], font=("Segoe UI", 9))
        self.subtitle_label.pack(side='left')

        # Theme toggle
        def _toggle_theme():
            self.theme.toggle()
            p = self.theme.palette

            self.configure(bg=p["bg"])
            tb = _shade(p["bg"], -0.25)
            self.topbar.configure(bg=tb)
            title_wrap.configure(bg=tb)
            self.title_label.configure(bg=tb, fg=p["text"])
            self.subtitle_label.configure(bg=tb, fg=p["muted"])

            self.toggle_btn.configure(bg=tb, fg=p["text"], activebackground=tb)

            self.paned.configure(bg=p["bg"])

            self.rules_panel.apply_palette(p)
            self.files_panel.apply_palette(p)
            self.folder_panel.apply_palette(p)
            self.preview_panel.apply_palette(p)

            
            if hasattr(self, "_editors"):
                for ed in list(self._editors):
                    try:
                        ed.apply_palette(p)
                    except Exception:
                        pass

        self.toggle_btn = tk.Button(
            self.topbar, text="🌙 / ☀️",
            command=_toggle_theme, bd=0,
            bg=self.topbar['bg'], fg=pal["text"],
            activebackground=self.topbar['bg']
        )
        self.toggle_btn.pack(side='right', padx=12, pady=8)                

        self.store = DataStore()  # loads rules/tags/folders

        # Window 3 columns
        self.paned = tk.PanedWindow(self, orient='horizontal', sashwidth=6, sashrelief='raised', bg=pal['bg'], bd=1)
        self.paned.pack(fill='both', expand=True)

        # Left - RulesPanel
        self.rules_panel = RulesPanel(self.paned, self.store)
        self.paned.add(self.rules_panel, minsize=300)

        # Middle - FileExplorerPanel
        self.files_panel = FileExplorerPanel(self.paned, self.store)
        self.paned.add(self.files_panel, minsize=700)

        self.files_panel.set_rules_provider(self.rules_panel.get_selected_rules)

        # Right -  Preview (top) + Folders (bottom)
        self.right_frame = tk.Frame(self.paned, bg=pal["bg"])
        self.paned.add(self.right_frame, minsize=260)

        self.right_frame.grid_rowconfigure(0, weight=3)   # preview area (bigger)
        self.right_frame.grid_rowconfigure(1, weight=2)   # folders area
        self.right_frame.grid_columnconfigure(0, weight=1)

        self.preview_panel = PreviewPanel(self.right_frame)
        self.preview_panel.grid(row=0, column=0, sticky="nsew")

        # Bottom: Folder panel 
        self.folder_panel = FolderPanel(self.right_frame, self.store,
                                        on_choose_for_files=self.files_panel.load_folder)
        self.folder_panel.grid(row=1, column=0, sticky="nsew")

        self.files_panel.set_rules_provider(self.rules_panel.get_selected_rules)

        self.files_panel.set_preview_panel(self.preview_panel)

        self.configure(bg=pal["bg"])
        self.rules_panel.apply_palette(pal)
        self.files_panel.apply_palette(pal)
        self.folder_panel.apply_palette(pal)
        self.preview_panel.apply_palette(pal)

        self.after(60, self._set_initial_sashes)

        default_path = self.store.default_folder_path()
        if default_path:
            self.files_panel.load_folder(default_path)

    def _set_initial_sashes(self):
        self.update_idletasks()
        total_w = self.paned.winfo_width()
        if total_w <= 1:
            self.after(60, self._set_initial_sashes)
            return
        left = int(total_w * 0.28)
        right = int(total_w * 0.82)
        try:
            self.paned.sashpos(0, left)
            self.paned.sashpos(1, right)
        except Exception:
            pass


if __name__ == '__main__':
    app = MainApp()
    app.mainloop()
